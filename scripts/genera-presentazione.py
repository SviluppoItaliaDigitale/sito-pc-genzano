#!/usr/bin/env python3
# Generatore del deck istituzionale di presentazione del sito PC Genzano.
# Output editabile, senza watermark, senza riferimenti a strumenti automatici/IA.
#
# REGOLA (CLAUDE.md § "Presentazione del sito (deck)"): rigenerare e ripubblicare
# il deck a OGNI sezione aggiunta/modificata/eliminata o aggiornamento di
# dati/contatti/fonti mostrati nelle slide. Output in static/manuali/.
#
# Uso:
#   python3 -m venv ~/.deck-venv && ~/.deck-venv/bin/pip install python-pptx pillow fonttools brotli
#   ~/.deck-venv/bin/python scripts/genera-presentazione.py   # + LibreOffice per il PDF
#
# Ricava icone (Bootstrap Icons self-hostate), logo e contatti dal repo stesso.
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
import os
import re
from PIL import Image, ImageDraw, ImageFont

REPO = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
_TTF = "/tmp/pcgenzano-bootstrap-icons.ttf"
if not os.path.exists(_TTF):
    from fontTools.ttLib import TTFont as _TTFont
    _bi = _TTFont(os.path.join(REPO, "static/vendor/bootstrap-icons/bootstrap-icons.woff2"))
    _bi.flavor = None
    _bi.save(_TTF)

# Titillium Web (font ufficiale AGID, lo stesso del portale): installa i TTF in
# ~/.fonts convertendoli dai woff2 del repo, così LibreOffice li incorpora come
# vettoriali nel PDF. Idempotente: salta se già presenti.
import subprocess
_FONTS_DIR = os.path.expanduser("~/.fonts")
_TW_SRC = os.path.join(REPO, "static/vendor/bootstrap-italia/fonts/Titillium_Web")
_tw_installed = False
for _w, _dst in (("regular", "TitilliumWeb-regular.ttf"), ("700", "TitilliumWeb-700.ttf"),
                 ("italic", "TitilliumWeb-italic.ttf")):
    _dstp = os.path.join(_FONTS_DIR, _dst)
    _srcp = os.path.join(_TW_SRC, "titillium-web-v10-latin-ext_latin-%s.woff2" % _w)
    if not os.path.exists(_dstp) and os.path.exists(_srcp):
        os.makedirs(_FONTS_DIR, exist_ok=True)
        from fontTools.ttLib import TTFont as _TWFont
        _t = _TWFont(_srcp); _t.flavor = None; _t.save(_dstp); _tw_installed = True
if _tw_installed:
    subprocess.run(["fc-cache", "-f", _FONTS_DIR], capture_output=True)

PRIMARY = RGBColor(0x00, 0x33, 0x66)   # blu istituzionale
PRIMARY_D = RGBColor(0x00, 0x1A, 0x33)
ACCENT  = RGBColor(0xFF, 0xBE, 0x2E)   # giallo accento
DARK    = RGBColor(0x1A, 0x1A, 0x1A)
GREY    = RGBColor(0x5A, 0x6B, 0x7B)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
# Titillium Web: font ufficiale AGID/Designers Italia, lo stesso del portale del
# Gruppo (coerenza di brand). Installato dal blocco sopra come TTF vettoriale,
# così LibreOffice non lo rasterizza nel PDF (il bug Calibri di maggio 2026).
FONT = "Titillium Web"
LOGO = os.path.join(REPO, "static/images/logo-pc-genzano.png")

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def logo_tr(slide, height_in=0.72, top_in=0.42, margin_in=0.5):
    pic = slide.shapes.add_picture(LOGO, 0, Inches(top_in), height=Inches(height_in))
    pic.left = SW - pic.width - Inches(margin_in)
    return pic


sec_targets = []  # (titolo capitolo, slide divisore, numero pagina)


def link_to_slide(run, from_slide, to_slide):
    rId = from_slide.part.relate_to(to_slide.part, RT.SLIDE)
    rPr = run._r.get_or_add_rPr()
    hlink = rPr.makeelement(qn('a:hlinkClick'),
                            {qn('r:id'): rId, 'action': 'ppaction://hlinksldjump'})
    rPr.append(hlink)


def add_shadow(shape):
    spPr = shape._element.spPr
    old = spPr.find(qn('a:effectLst'))
    if old is not None:
        spPr.remove(old)
    el = spPr.makeelement(qn('a:effectLst'), {})
    sh = el.makeelement(qn('a:outerShdw'),
                        {'blurRad': '69850', 'dist': '34290',
                         'dir': '5400000', 'rotWithShape': '0'})
    clr = sh.makeelement(qn('a:srgbClr'), {'val': '001A33'})
    a = clr.makeelement(qn('a:alpha'), {'val': '46000'})
    clr.append(a); sh.append(clr); el.append(sh)
    spPr.append(el)


def _hex(c):
    return '%02X%02X%02X' % (c[0], c[1], c[2])


def _darken(c, f=0.7):
    return RGBColor(int(c[0] * f), int(c[1] * f), int(c[2] * f))


def _lighten(c, f=1.15):
    return RGBColor(min(255, int(c[0] * f)), min(255, int(c[1] * f)), min(255, int(c[2] * f)))


def grad_fill(shape, c1, c2, ang=90):
    """Sostituisce il riempimento dello shape con un gradiente lineare c1→c2.
    ang in gradi (90 = verticale dall'alto, 0 = orizzontale da sinistra)."""
    spPr = shape._element.spPr
    for tag in ('a:solidFill', 'a:noFill', 'a:gradFill', 'a:pattFill', 'a:blipFill'):
        e = spPr.find(qn(tag))
        if e is not None:
            spPr.remove(e)
    g = spPr.makeelement(qn('a:gradFill'), {})
    lst = g.makeelement(qn('a:gsLst'), {})
    for pos, col in ((0, c1), (100000, c2)):
        gs = g.makeelement(qn('a:gs'), {'pos': str(pos)})
        clr = gs.makeelement(qn('a:srgbClr'), {'val': _hex(col)})
        gs.append(clr)
        lst.append(gs)
    g.append(lst)
    g.append(g.makeelement(qn('a:lin'), {'ang': str(int(ang * 60000)), 'scaled': '1'}))
    ln = spPr.find(qn('a:ln'))
    if ln is not None:
        ln.addprevious(g)
    else:
        spPr.append(g)
    return shape


def grad_rect(slide, l, t, w, h, c1, c2, ang=90):
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    sp.fill.solid(); sp.fill.fore_color.rgb = c1
    sp.line.fill.background(); sp.shadow.inherit = False
    grad_fill(sp, c1, c2, ang)
    return sp


def _alpha(shape, pct):
    clr = shape._element.spPr.find(qn('a:solidFill')).find(qn('a:srgbClr'))
    clr.append(clr.makeelement(qn('a:alpha'), {'val': str(int(pct * 1000))}))
    return shape


def _norm(slide):
    # rimuove eventuali placeholder vuoti
    for ph in list(slide.placeholders):
        ph._element.getparent().remove(ph._element)


def rect(slide, l, t, w, h, color, line=False):
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    sp.fill.solid(); sp.fill.fore_color.rgb = color
    if line:
        sp.line.color.rgb = color
    else:
        sp.line.fill.background()
    sp.shadow.inherit = False
    return sp


def textbox(slide, l, t, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.05); tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02); tf.margin_bottom = Inches(0.02)
    return tf


def setrun(r, text, size, color, bold=False, italic=False):
    r.text = text
    r.font.size = Pt(size); r.font.color.rgb = color
    r.font.bold = bold; r.font.italic = italic; r.font.name = FONT


def footer(slide, n, dark=False, nav=True):
    line_col = RGBColor(0x33, 0x4E, 0x66) if dark else RGBColor(0xD8, 0xDE, 0xE4)
    txt_col = RGBColor(0x9D, 0xB0, 0xC2) if dark else GREY
    rect(slide, Inches(0.6), Inches(7.02), Inches(12.13), Pt(1.2), line_col)
    tf = textbox(slide, Inches(0.6), Inches(7.06), Inches(4.5), Inches(0.35))
    setrun(tf.paragraphs[0].add_run(),
           "Gruppo Comunale Volontari di PC · Genzano di Roma", 9, txt_col)
    tf2 = textbox(slide, Inches(11.7), Inches(7.06), Inches(1.0), Inches(0.35))
    p2 = tf2.paragraphs[0]; p2.alignment = PP_ALIGN.RIGHT
    setrun(p2.add_run(), str(n), 9, txt_col, bold=True)
    ag = globals().get('agenda')
    if nav and ag is not None and slide is not ag:
        tfn = textbox(slide, Inches(5.17), Inches(7.03), Inches(3.0), Inches(0.36))
        pn = tfn.paragraphs[0]; pn.alignment = PP_ALIGN.CENTER
        r = pn.add_run(); setrun(r, "« Torna all'agenda", 9.5, (ACCENT if dark else PRIMARY), bold=True)
        link_to_slide(r, slide, ag)


def cover(title, subtitle, eyebrow="PROTEZIONE CIVILE — GENZANO DI ROMA"):
    s = prs.slides.add_slide(BLANK); _norm(s)
    grad_rect(s, 0, 0, SW, SH, PRIMARY_D, PRIMARY, ang=115)
    _alpha(shp(s, MSO_SHAPE.OVAL, Inches(8.7), Inches(3.4), Inches(7.2), Inches(7.2), BLUE2), 14)
    _alpha(shp(s, MSO_SHAPE.OVAL, Inches(10.7), Inches(-1.8), Inches(4.2), Inches(4.2), ACCENT), 9)
    rect(s, 0, 0, Inches(0.22), SH, ACCENT)
    s.shapes.add_picture(LOGO, Inches(1.0), Inches(0.6), height=Inches(1.25))
    tf = textbox(s, Inches(1.0), Inches(2.0), Inches(11.3), Inches(0.6))
    setrun(tf.paragraphs[0].add_run(), eyebrow, 14, ACCENT, bold=True)
    tf2 = textbox(s, Inches(1.0), Inches(2.6), Inches(11.3), Inches(2.0))
    setrun(tf2.paragraphs[0].add_run(), title, 44, WHITE, bold=True)
    rect(s, Inches(1.05), Inches(4.7), Inches(3.2), Pt(3), ACCENT)
    tf3 = textbox(s, Inches(1.0), Inches(4.95), Inches(11.3), Inches(1.0))
    setrun(tf3.paragraphs[0].add_run(), subtitle, 20, RGBColor(0xCF, 0xDD, 0xEC))
    return s


def divider(chapter, num):
    s = prs.slides.add_slide(BLANK); _norm(s)
    grad_rect(s, 0, 0, SW, SH, RGBColor(0x00, 0x10, 0x21), PRIMARY, ang=120)
    _alpha(shp(s, MSO_SHAPE.OVAL, Inches(-2.0), Inches(3.7), Inches(6.0), Inches(6.0), BLUE2), 16)
    logo_tr(s, height_in=0.95, top_in=0.55)
    rect(s, Inches(1.0), Inches(3.0), Inches(2.4), Pt(4), ACCENT)
    # Titolo confinato alla metà sinistra: il menu sezioni occupa la destra (x>=7).
    tf = textbox(s, Inches(1.0), Inches(3.25), Inches(5.5), Inches(2.7))
    setrun(tf.paragraphs[0].add_run(), chapter, 32, WHITE, bold=True)
    tf2 = textbox(s, Inches(1.0), Inches(2.4), Inches(5.5), Inches(0.5))
    setrun(tf2.paragraphs[0].add_run(), "SEZIONE", 13, ACCENT, bold=True)
    footer(s, num, dark=True)
    sec_targets.append((chapter, s, num))
    return s


def content(title, bullets, num, fonti=None):
    # Ogni punto diventa una CARD numerata; il layout si adatta al numero di punti.
    s = base_slide(title, num)
    ic = icon_png(_icon_for(title))
    n = len(bullets)
    top = 2.15
    bottom = 6.25 if fonti else 6.7
    if n >= 5:                              # pila verticale
        gap = 0.2; avail = bottom - top
        h = min(1.0, (avail - gap * (n - 1)) / n)
        used = h * n + gap * (n - 1); y0 = top + max(0, (avail - used) / 2)
        for i, b in enumerate(bullets):
            _card(s, 0.85, y0 + i * (h + gap), 11.6, h, i, b, 'v', icon=ic)
    elif n == 4:                            # griglia 2x2
        w = 5.71; gap = 0.18; vgap = 0.34; avail = bottom - top
        h = min(2.02, (avail - vgap) / 2); y0 = top + (avail - (h * 2 + vgap)) / 2
        for i, b in enumerate(bullets):
            _card(s, 0.85 + (i % 2) * (w + gap), y0 + (i // 2) * (h + vgap), w, h, i, b, 'g', icon=ic)
    else:                                   # 1-3 colonne affiancate, centrate
        gap = 0.2; w = (11.6 - (n - 1) * gap) / n
        h = min(2.95, bottom - top); y = top + (bottom - top - h) / 2
        for i, b in enumerate(bullets):
            _card(s, 0.85 + i * (w + gap), y, w, h, i, b, 'g', icon=ic)
    if fonti:
        fb = textbox(s, Inches(0.85), Inches(6.35), Inches(11.6), Inches(0.45))
        p = fb.paragraphs[0]
        setrun(p.add_run(), "Fonti e standard: ", 12, PRIMARY, bold=True)
        setrun(p.add_run(), fonti, 12, GREY, italic=True)
    return s


# ── Costruttori di slide-diagramma (infografiche vettoriali editabili) ──
def base_slide(title, num):
    s = prs.slides.add_slide(BLANK); _norm(s)
    rect(s, 0, 0, SW, SH, WHITE)
    hb = grad_rect(s, Inches(0.22), 0, SW - Inches(0.22), Inches(1.74), WHITE, RGBColor(0xEC, 0xF1, 0xF8), ang=90)
    add_shadow(hb)
    rect(s, Inches(0.75), Inches(1.74), Inches(11.83), Pt(0.9), RGBColor(0xDD, 0xE5, 0xEC))
    grad_rect(s, 0, 0, Inches(0.22), SH, PRIMARY, BLUE2, ang=90)
    logo_tr(s, height_in=0.72, top_in=0.45)
    chip = shp(s, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(0.52), Inches(0.92), Inches(0.92), PRIMARY, shadow=True)
    grad_fill(chip, BLUE2, PRIMARY_D, ang=120)
    _ci = icon_png(_icon_for(title))
    if _ci:
        s.shapes.add_picture(_ci, Inches(0.94), Inches(0.76), height=Inches(0.44))
    tf = textbox(s, Inches(1.82), Inches(0.5), Inches(9.5), Inches(0.98), anchor=MSO_ANCHOR.MIDDLE)
    setrun(tf.paragraphs[0].add_run(), title, 27, PRIMARY, bold=True)
    grad_rect(s, Inches(1.82), Inches(1.55), Inches(2.2), Pt(3.2), ACCENT, RGBColor(0xE0, 0x90, 0x10), ang=0)
    footer(s, num)
    return s


def shp(slide, st, x, y, w, h, fill, line=None, lw=1.0, shadow=False):
    o = slide.shapes.add_shape(st, x, y, w, h)
    o.fill.solid(); o.fill.fore_color.rgb = fill
    if line is None:
        o.line.fill.background()
    else:
        o.line.color.rgb = line; o.line.width = Pt(lw)
    o.shadow.inherit = False
    if shadow:
        add_shadow(o)
    tf = o.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.12); tf.margin_right = Inches(0.12)
    tf.margin_top = Inches(0.04); tf.margin_bottom = Inches(0.04)
    return o


def put(shape, runs, align=PP_ALIGN.CENTER):
    tf = shape.text_frame
    first = True
    for (txt, sz, col, bd) in runs:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align; p.line_spacing = 1.0
        setrun(p.add_run(), txt, sz, col, bold=bd)


GREEN = RGBColor(0x1E, 0x84, 0x49); YELLOW = RGBColor(0xF1, 0xC4, 0x0F)
ORANGE = RGBColor(0xE6, 0x7E, 0x22); RED = RGBColor(0xC0, 0x39, 0x2B)
GREENC = RGBColor(0x15, 0x80, 0x3D); BLUE2 = RGBColor(0x07, 0x59, 0x85)
PALETTE = [PRIMARY, GREENC, RGBColor(0xB4, 0x53, 0x09), BLUE2,
           RGBColor(0x7C, 0x3A, 0xED), RGBColor(0x0E, 0x74, 0x90)]


# ── Icone (Bootstrap Icons self-hostate → PNG bianchi/colorati) ──
_ICON_DIR = "/tmp/deck_icons"
os.makedirs(_ICON_DIR, exist_ok=True)
_CSS_BI = os.path.join(REPO, "static/vendor/bootstrap-icons/bootstrap-icons.min.css")
_CPMAP = dict(re.findall(r'\.bi-([a-z0-9-]+)::before\{content:"\\([0-9a-fA-F]+)"\}',
                         open(_CSS_BI, encoding="utf-8", errors="replace").read()))
_IFONT = ImageFont.truetype(_TTF, 300)


def icon_png(name, rgb=(255, 255, 255)):
    if not name or name not in _CPMAP:
        return None
    tag = "%s_%02x%02x%02x" % (name, rgb[0], rgb[1], rgb[2])
    path = os.path.join(_ICON_DIR, tag + ".png")
    if not os.path.exists(path):
        ch = chr(int(_CPMAP[name], 16))
        img = Image.new("RGBA", (360, 360), (0, 0, 0, 0))
        d = ImageDraw.Draw(img)
        bb = d.textbbox((0, 0), ch, font=_IFONT)
        w0, h0 = bb[2] - bb[0], bb[3] - bb[1]
        d.text(((360 - w0) / 2 - bb[0], (360 - h0) / 2 - bb[1]), ch, font=_IFONT, fill=rgb + (255,))
        img.save(path)
    return path


def _icon_for(title):
    t = title.lower()
    table = [
        ("rischi e prevenzione", "exclamation-triangle-fill"),
        ("estensione digitale", "diagram-3-fill"),
        ("allerte", "cloud-lightning-rain-fill"),
        ("piattaforma", "people-fill"), ("principi", "compass"),
        ("conformità", "patch-check-fill"), ("architettura", "diagram-3-fill"),
        ("cosa fare", "signpost-2-fill"), ("modalità emergenza", "exclamation-octagon-fill"),
        ("cruscotto", "database-fill-gear"), ("cartografia", "geo-alt-fill"),
        ("numeri utili", "telephone-fill"), ("piano familiare", "house-heart-fill"),
        ("kit pronti", "backpack-fill"), ("storia", "clock-history"),
        ("kit didattici", "mortarboard-fill"), ("percorsi didattici", "journal-bookmark-fill"),
        ("schede", "printer-fill"), ("storie e racconti", "book-half"),
        ("giochi", "controller"), ("wcag", "universal-access"),
        ("lettura facilitata", "volume-up-fill"), ("pittogrammi", "image-fill"),
        ("italiano semplice", "chat-left-text-fill"), ("lingua dei segni", "hand-index-thumb-fill"),
        ("braille", "grid-3x3-gap-fill"), ("più lingue", "translate"),
        ("diventa volontario", "person-plus-fill"), ("faq", "question-circle-fill"),
        ("normativa", "journal-text"), ("area download", "download"),
        ("strumenti", "speedometer2"), ("audio e podcast", "headphones"),
        ("open data", "bar-chart-fill"), ("pagine legali", "file-earmark-text-fill"),
        ("comunicazioni e notizie", "megaphone-fill"), ("comunicazione di crisi", "broadcast"),
        ("canali", "rss-fill"), ("infrastruttura", "hdd-network-fill"),
        ("fonti e dati", "bookmark-check-fill"), ("punti di forza", "stars"),
    ]
    for k, ic in table:
        if k in t:
            return ic
    return "info-circle-fill"


def _card(s, xi, yi, wi, hi, idx, text, mode, icon=None):
    x, y, w, h = Inches(xi), Inches(yi), Inches(wi), Inches(hi)
    c = PALETTE[idx % len(PALETTE)]
    grad_fill(shp(s, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h, RGBColor(0xF4, 0xF7, 0xFB),
                  line=RGBColor(0xD8, 0xDE, 0xE4), shadow=True),
              RGBColor(0xFF, 0xFF, 0xFF), RGBColor(0xEA, 0xF0, 0xF7), ang=90)
    if mode == 'v':
        bsz, bcx, bcy = 0.56, xi + 0.6, yi + hi / 2
        tb = textbox(s, Inches(xi + 1.18), y, Inches(wi - 1.35), h, anchor=MSO_ANCHOR.MIDDLE)
        p = tb.paragraphs[0]; p.alignment = PP_ALIGN.LEFT; p.line_spacing = 1.1
        setrun(p.add_run(), text, 15.5, DARK)
    else:
        bsz, bcx, bcy = 0.66, xi + wi / 2, yi + hi * 0.36
        ty = yi + hi * 0.36 + bsz / 2 + 0.08
        tb = textbox(s, Inches(xi + 0.25), Inches(ty), Inches(wi - 0.5),
                     Inches(yi + hi - 0.15 - ty), anchor=MSO_ANCHOR.MIDDLE)
        p = tb.paragraphs[0]; p.alignment = PP_ALIGN.CENTER; p.line_spacing = 1.18
        setrun(p.add_run(), text, 15, DARK)
    badge = shp(s, MSO_SHAPE.OVAL, Inches(bcx - bsz / 2), Inches(bcy - bsz / 2),
                Inches(bsz), Inches(bsz), c)
    grad_fill(badge, c, _darken(c, 0.6), ang=90)
    if icon:
        isz = bsz * 0.56
        s.shapes.add_picture(icon, Inches(bcx - isz / 2), Inches(bcy - isz / 2), height=Inches(isz))
    else:
        put(badge, [(str(idx + 1), 17, WHITE, True)])


def slide_funzioni(num):
    s = base_slide("Un'estensione digitale del Servizio Nazionale", num)
    it = textbox(s, Inches(0.75), Inches(1.85), Inches(11.9), Inches(0.6))
    setrun(it.paragraphs[0].add_run(),
           "Il Gruppo opera nelle quattro funzioni del Servizio Nazionale della Protezione Civile (D.Lgs. 1/2018).",
           15, DARK)
    cards = [("Previsione", "Scenari e monitoraggio del rischio.", PRIMARY, "binoculars-fill", (0x00, 0x33, 0x66)),
             ("Prevenzione", "Riduzione del rischio e autoprotezione.", GREENC, "shield-check", (0x15, 0x80, 0x3D)),
             ("Soccorso", "Supporto alla popolazione durante l'evento.", RGBColor(0xB4, 0x53, 0x09), "heart-pulse-fill", (0xB4, 0x53, 0x09)),
             ("Superamento", "Assistenza e ripristino dopo l'evento.", BLUE2, "arrow-repeat", (0x07, 0x59, 0x85))]
    x0, w, gap, y, h = 0.85, 2.85, 0.18, 2.8, 3.05
    for i, (t, d, c, icn, rgb) in enumerate(cards):
        bxf = x0 + i * (w + gap); bx = Inches(bxf)
        grad_fill(shp(s, MSO_SHAPE.ROUNDED_RECTANGLE, bx, Inches(y), Inches(w), Inches(h),
                      RGBColor(0xF4, 0xF7, 0xFB), line=RGBColor(0xD8, 0xDE, 0xE4), shadow=True),
                  RGBColor(0xFF, 0xFF, 0xFF), RGBColor(0xEA, 0xF0, 0xF7), ang=90)
        grad_fill(shp(s, MSO_SHAPE.RECTANGLE, bx, Inches(y), Inches(w), Inches(0.14), c),
                  _lighten(c, 1.1), _darken(c, 0.85), ang=0)
        ip = icon_png(icn, rgb=rgb)
        if ip:
            s.shapes.add_picture(ip, Inches(bxf + w / 2 - 0.42), Inches(y + 0.45), height=Inches(0.84))
        body = textbox(s, bx, Inches(y + 1.5), Inches(w), Inches(h - 1.65))
        bp = body.paragraphs[0]; bp.alignment = PP_ALIGN.CENTER
        setrun(bp.add_run(), t, 17, c, bold=True)
        d2 = body.add_paragraph(); d2.alignment = PP_ALIGN.CENTER
        d2.space_before = Pt(8); d2.line_spacing = 1.1
        setrun(d2.add_run(), d, 12.5, DARK)
    fn = textbox(s, Inches(0.85), Inches(6.25), Inches(11.6), Inches(0.4))
    setrun(fn.paragraphs[0].add_run(),
           "Riferimenti: D.Lgs. 1/2018 · Direttiva PCM 30 aprile 2021.", 11, GREY, italic=True)
    return s


def slide_sistema(num):
    s = base_slide("Chi siamo nel sistema di Protezione Civile", num)
    _BASE = os.path.join(REPO, "static", "images") + os.sep
    items = [("Servizio Nazionale della Protezione Civile", "D.Lgs. 1/2018", PRIMARY, None),
             ("Coordinamento Regionale FE.PI.VOL.",
              "Federazione Pronto Intervento Volontariato ODV — volontariato di PC del Lazio",
              BLUE2, _BASE + "logo-fepivol.png"),
             ("Gruppo Comunale Volontari di Protezione Civile di Genzano di Roma",
              "Sede operativa: Via Sicilia 13-15", GREENC, _BASE + "logo-pc-genzano.png")]
    w, h, gap, y = 7.6, 1.18, 0.5, 2.2
    xf = (13.333 - w) / 2; x = Inches(xf)
    for i, (t, sub, c, logo) in enumerate(items):
        yf = y + i * (h + gap); yy = Inches(yf)
        b = shp(s, MSO_SHAPE.ROUNDED_RECTANGLE, x, yy, Inches(w), Inches(h), c, shadow=True)
        put(b, [(t, 15, WHITE, True), (sub, 10.5, RGBColor(0xCF, 0xDD, 0xEC), False)])
        if logo:
            s.shapes.add_picture(logo, Inches(1.25), Inches(yf + h / 2 - 0.52), height=Inches(1.04))
        if i < len(items) - 1:
            shp(s, MSO_SHAPE.DOWN_ARROW, Inches(6.4665), Inches(yf + h + 0.03),
                Inches(0.4), Inches(0.4), ACCENT)
    return s


def slide_cruscotto(num):
    s = base_slide("Cruscotto del territorio — dati in tempo reale", num)
    srcs = [("Scientifiche", "INGV (sismica, vulcani)"),
            ("Meteo e satellite", "ItaliaMeteo, EUMETSAT, ECMWF, Copernicus"),
            ("Istituzionali", "DPC, ARPA Lazio")]
    sx, sw, shh, sy0, sgap = Inches(0.85), Inches(3.5), 0.95, 2.35, 0.4
    for i, (t, d) in enumerate(srcs):
        yy = Inches(sy0 + i * (shh + sgap))
        b = shp(s, MSO_SHAPE.ROUNDED_RECTANGLE, sx, yy, sw, Inches(shh),
                RGBColor(0xF4, 0xF7, 0xFB), line=RGBColor(0xCB, 0xD7, 0xE6))
        put(b, [(t, 13, PRIMARY, True), (d, 10.5, DARK, False)])
        shp(s, MSO_SHAPE.RIGHT_ARROW, Inches(4.5), Inches(sy0 + i * (shh + sgap) + shh / 2 - 0.18),
            Inches(0.62), Inches(0.36), ACCENT)
    cb = shp(s, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.32), Inches(2.7), Inches(3.05), Inches(2.6),
             PRIMARY, shadow=True)
    grad_fill(cb, BLUE2, PRIMARY_D, ang=120)
    put(cb, [("Cruscotto del territorio", 15, WHITE, True),
             ("Aggregazione privacy-first", 11, RGBColor(0xCF, 0xDD, 0xEC), False),
             ("Dati aperti (CC BY 4.0)", 11, RGBColor(0xCF, 0xDD, 0xEC), False)])
    shp(s, MSO_SHAPE.RIGHT_ARROW, Inches(8.5), Inches(3.85), Inches(0.62), Inches(0.36), ACCENT)
    al = [("Verde — nessuna allerta", GREEN, WHITE), ("Giallo — ordinaria", YELLOW, PRIMARY_D),
          ("Arancione — elevata", ORANGE, WHITE), ("Rosso — massima", RED, WHITE)]
    ax, aw, ah, ay0, agp = Inches(9.35), Inches(3.3), 0.62, 2.4, 0.22
    for i, (t, c, tc) in enumerate(al):
        b = shp(s, MSO_SHAPE.ROUNDED_RECTANGLE, ax, Inches(ay0 + i * (ah + agp)), aw, Inches(ah), c, shadow=True)
        grad_fill(b, _lighten(c, 1.12), _darken(c, 0.82), ang=90)
        put(b, [(t, 12, tc, True)])
    fn = textbox(s, Inches(0.85), Inches(6.25), Inches(11.6), Inches(0.5))
    setrun(fn.paragraphs[0].add_run(),
           "Fonti: INGV · DPC · ItaliaMeteo · CIMA · EUMETSAT · ECMWF · Copernicus · ARPA Lazio. Codici colore: Centro Funzionale Regionale del Lazio (ISO 22324).",
           11, GREY, italic=True)
    return s


def slide_allerte(num):
    s = base_slide("Allerte meteo e codici colore", num)
    it = textbox(s, Inches(0.75), Inches(1.85), Inches(11.9), Inches(0.6))
    setrun(it.paragraphs[0].add_run(),
           "Quattro livelli ufficiali del Centro Funzionale Regionale del Lazio. La previsione (allerta) è distinta dall'evento in corso (emergenza).",
           14, DARK)
    bars = [("VERDE", "Nessuna allerta — situazione ordinaria", GREEN, WHITE),
            ("GIALLO", "Allerta ordinaria — fenomeni di intensità moderata", YELLOW, PRIMARY_D),
            ("ARANCIONE", "Allerta elevata — fenomeni intensi, effetti rilevanti", ORANGE, WHITE),
            ("ROSSO", "Allerta massima — fenomeni molto intensi, rischio elevato", RED, WHITE)]
    x, w, h, y0, gp = Inches(0.85), Inches(11.6), 0.78, 2.55, 0.2
    for i, (lv, d, c, tc) in enumerate(bars):
        b = shp(s, MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(y0 + i * (h + gp)), w, Inches(h), c, shadow=True)
        grad_fill(b, _lighten(c, 1.12), _darken(c, 0.82), ang=90)
        tf = b.text_frame; tf.margin_left = Inches(0.3)
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
        setrun(p.add_run(), lv + "    ", 16, tc, bold=True)
        setrun(p.add_run(), d, 13, tc, bold=False)
    fn = textbox(s, Inches(0.85), Inches(6.5), Inches(11.6), Inches(0.4))
    setrun(fn.paragraphs[0].add_run(),
           "Fonte: Centro Funzionale Regionale del Lazio · standard ISO 22324.", 11, GREY, italic=True)
    return s


def slide_rischi(num):
    s = base_slide("Rischi e prevenzione — copertura del territorio", num)
    it = textbox(s, Inches(0.75), Inches(1.85), Inches(11.9), Inches(0.7))
    setrun(it.paragraphs[0].add_run(),
           "Il sito tratta sistematicamente i rischi dei Castelli Romani. Ogni scheda segue la stessa struttura: prima, durante, dopo, cosa non fare, chi chiamare.",
           14, DARK)
    risks = [("Sismico", RGBColor(0xB4, 0x53, 0x09)), ("Idrogeologico", RGBColor(0x03, 0x69, 0xA1)),
             ("Incendi boschivi", RGBColor(0xC1, 0x12, 0x1F)), ("Ondate di calore", RGBColor(0xEA, 0x58, 0x0C)),
             ("Vento forte", RGBColor(0x0E, 0x74, 0x90)), ("Temporali intensi", RGBColor(0x43, 0x38, 0xCA)),
             ("Blackout", RGBColor(0x47, 0x55, 0x69)), ("Vulcanico (Colli Albani)", RGBColor(0x7F, 0x1D, 0x1D))]
    x0, w, gap, y0, h, vg = 0.85, 2.85, 0.18, 2.85, 1.45, 0.3
    for i, (t, c) in enumerate(risks):
        col, row = i % 4, i // 4
        b = shp(s, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x0 + col * (w + gap)),
                Inches(y0 + row * (h + vg)), Inches(w), Inches(h), c, shadow=True)
        grad_fill(b, _lighten(c, 1.16), _darken(c, 0.74), ang=90)
        put(b, [(t, 15, WHITE, True)])
    fn = textbox(s, Inches(0.85), Inches(6.35), Inches(11.6), Inches(0.4))
    setrun(fn.paragraphs[0].add_run(),
           "Fonti: DPC «Io non rischio» · INGV · ISPRA · CNR-IRPI · ReLUIS.", 11, GREY, italic=True)
    return s


# ─────────────────────────── CONTENUTI ───────────────────────────
n = 0
def N():
    global n; n += 1; return n

cover("Il sito istituzionale:\nstruttura e contenuti",
      "Gruppo Comunale Volontari di Protezione Civile di Genzano di Roma")
n = 1

# Agenda (slide 2) — il corpo con i link cliccabili è popolato in fondo,
# quando tutte le slide-divisore di destinazione esistono già.
agenda = prs.slides.add_slide(BLANK); _norm(agenda)
rect(agenda, 0, 0, SW, SH, WHITE)
rect(agenda, 0, 0, Inches(0.22), SH, PRIMARY)
logo_tr(agenda, height_in=0.72, top_in=0.45)
_tfh = textbox(agenda, Inches(0.7), Inches(0.55), Inches(10.6), Inches(1.1))
setrun(_tfh.paragraphs[0].add_run(), "Agenda", 28, PRIMARY, bold=True)
rect(agenda, Inches(0.75), Inches(1.55), Inches(2.2), Pt(3), ACCENT)
_sub = textbox(agenda, Inches(0.85), Inches(6.45), Inches(11.6), Inches(0.4))
setrun(_sub.paragraphs[0].add_run(),
       "Le voci sono cliccabili: portano alla sezione corrispondente. Da ogni slide il piè di pagina riporta qui con «Torna all'agenda».",
       12, GREY, italic=True)
footer(agenda, 2, nav=False)
n = 2

slide_funzioni(N())

content("Una piattaforma per l'intera comunità", [
    "Famiglie e cittadini; anziani e persone fragili.",
    "Persone con disabilità sensoriali, cognitive e motorie.",
    "Scuole e studenti; volontari e operatori.",
    "Progettazione mobile-first: consultabile in movimento e in condizioni critiche.",
], N())

content("Principi guida", [
    "Affidabilità e correttezza istituzionale prima di tutto.",
    "Accessibilità universale, con le WCAG 2.2 AA come soglia minima.",
    "Chiarezza per il cittadino e sobrietà del linguaggio.",
    "Ogni dato è verificato e accompagnato dalla sua fonte.",
], N())

content("Conformità e standard di riferimento", [
    "Linee guida di design AgID e design system Bootstrap Italia / Designers Italia.",
    "Accessibilità WCAG 2.2 AA; Legge Stanca (L. 4/2004) e D.Lgs. 106/2018.",
    "Codice della Protezione Civile e direttive nazionali.",
    "Standard internazionali ISO per la gestione delle emergenze e del rischio.",
], N(), fonti="AgID · W3C · DPC · ISO/TC 292")

content("Architettura dell'informazione", [
    "Sette aree principali: Per il Cittadino, Per le Scuole, Accessibilità e Supporti.",
    "Volontariato, Risorse, Comunicazioni, Contatti.",
    "Struttura piatta: ogni risorsa critica raggiungibile in pochi clic.",
    "Ricerca interna, mappa del sito e codici QR per l'accesso rapido.",
], N())

divider("Per il Cittadino", N())

content("Cosa fare adesso e assistente guidato", [
    "Pagina d'ingresso rapida alle azioni di autoprotezione.",
    "Assistente guidato a domande e risposte predefinite, lungo un percorso a scelte.",
    "Indirizza al comportamento corretto in pochi passaggi.",
    "Sempre il richiamo al 112 per le emergenze reali.",
], N())

slide_allerte(N())

content("Modalità emergenza del sito", [
    "In emergenza il sito mostra un avviso ben visibile su ogni pagina.",
    "Pagina di emergenza ultraleggera per rete debole o congestionata.",
    "Informazioni essenziali e numeri utili messi in primo piano.",
    "Progettata per i momenti di reale criticità.",
], N())

slide_cruscotto(N())

content("Cruscotto — le fonti dei dati", [
    "INGV per la sismicità e i vulcani; Dipartimento della Protezione Civile per il radar pioggia.",
    "ItaliaMeteo (modello ICON-2I) ed EUMETSAT per previsioni e satellite.",
    "EFFIS / Copernicus per gli incendi; ARPA Lazio per la qualità dell'aria regionale; Copernicus CAMS per la qualità dell'aria su scala europea.",
    "Copernicus EMS Rapid Mapping per le attivazioni di mappatura emergenze in Europa.",
    "Rischio sismico comunale dalla piattaforma SicuroPiù.",
], N(), fonti="INGV · DPC · ItaliaMeteo · EUMETSAT · Copernicus (CAMS · EMS · EFFIS) · ARPA Lazio · EUCENTRE/ReLUIS")

slide_rischi(N())

content("Cartografia e aree di emergenza", [
    "Mappa delle aree di attesa e di emergenza comunali.",
    "Riferimento diretto al Piano di Emergenza Comunale.",
    "Orientamento immediato per il cittadino.",
], N())

content("Numeri utili", [
    "Numero Unico di Emergenza 112 sempre in evidenza.",
    "Numeri di servizio per le segnalazioni non urgenti.",
    "Indicazioni chiare su quando e chi chiamare.",
], N())

content("Piano familiare di emergenza", [
    "Guida passo-passo per preparare il nucleo domestico.",
    "Versione compilabile e stampabile.",
    "Punti di ritrovo, contatti e kit essenziale.",
], N())

content("Kit pronti per situazioni vulnerabili", [
    "Schede operative per le categorie fragili: anziani, bambini, persone con disabilità, caregiver e altre.",
    "Materiali stampabili pronti all'uso.",
    "Riferimenti a standard internazionali di assistenza umanitaria.",
], N(), fonti="IFRC · WHO · Sphere Handbook")

content("Storia del territorio", [
    "Cronologia degli eventi rilevanti dei Castelli Romani.",
    "Ogni voce è corredata da fonte istituzionale.",
    "La memoria come strumento di consapevolezza del rischio.",
], N())

divider("Per le Scuole", N())

content("Kit didattici per fascia scolastica", [
    "Materiali per infanzia, primaria, secondaria di primo e di secondo grado.",
    "Linguaggio e attività calibrati per ogni età.",
    "Coerenti con l'insegnamento dell'Educazione Civica.",
], N())

content("Percorsi didattici ed Educazione Civica", [
    "Percorsi pronti all'uso per i docenti.",
    "Collegamento esplicito all'Educazione Civica.",
    "Obiettivo: cultura della prevenzione fin da piccoli.",
], N())

content("Schede didattiche stampabili", [
    "Schede fotocopiabili per le attività in classe.",
    "Organizzate per livello scolastico.",
    "Pronte all'uso, scaricabili e riproducibili liberamente.",
], N())

content("Storie e racconti", [
    "Narrazioni illustrate sulla sicurezza per i più piccoli.",
    "Disponibili in lettura accessibile e in ascolto.",
    "Avvicinano i bambini ai comportamenti corretti con il racconto.",
], N())

content("Giochi della sicurezza", [
    "Giochi educativi divisi per fascia d'età.",
    "I comportamenti corretti si imparano attraverso il gioco.",
    "Supporto e teoria di rinforzo integrati.",
], N())

content("Catalogo informativo dei giochi", [
    "Una scheda dedicata per ogni gioco, divisa per fascia d'età.",
    "Per ogni gioco: obiettivo didattico, durata, accessibilità e versione interattiva.",
    "Pensato per docenti, genitori, motori di ricerca e lettori di schermo.",
], N(), fonti="D.M. 183/2024 (Educazione Civica) · WCAG 2.2 AA")

content("Laboratorio meteo — esplorare i dati del clima", [
    "Strumento interattivo: si scelgono luogo, periodo e variabile e si costruisce il grafico.",
    "Dati storici reali dei Castelli Romani (rianalisi ERA5) ed esempi già pronti sul clima locale.",
    "Ogni grafico ha la tabella dati, la stampa, il download (CSV e immagine) e un link condivisibile.",
], N(), fonti="ERA5 · Open-Meteo (CC BY 4.0) · Copernicus Climate Data Store")

content("Dossier interattivi — racconti visivi", [
    "Storie che si leggono scorrendo: immagini a tutto schermo, dati animati e mappe da esplorare.",
    "Qualità da rivista digitale, ma interamente accessibile (WCAG 2.2 AA) e senza servizi di terzi.",
    "Ogni pagina è condivisibile; il primo dossier racconta la Terra vista dai satelliti.",
], N(), fonti="Immagini con licenza aperta: NASA · Copernicus / ESA")

divider("Accessibilità e inclusione", N())

content("Conformità WCAG 2.2 AA e barra strumenti", [
    "Sito conforme alle WCAG 2.2 di livello AA.",
    "Barra strumenti di lettura: dimensione del testo, contrasti, caratteri ad alta leggibilità.",
    "Pausa delle animazioni, spaziatura ampia, modalità lettura.",
], N(), fonti="W3C WCAG 2.2 · Legge Stanca (L. 4/2004)")

content("Lettura facilitata", [
    "Lettura ad alta voce dei contenuti.",
    "Tempo di lettura stimato e sillabazione automatica.",
    "Glossario esplicativo a comparsa sui termini tecnici e le sigle.",
], N())

content("Pittogrammi standardizzati", [
    "Simboli ISO 7010 per la segnaletica di sicurezza.",
    "Pittogrammi ARASAAC per la comprensione cognitiva.",
    "Supporto a bambini, anziani e persone con difficoltà di lettura.",
], N(), fonti="ISO 7010 · ARASAAC")

content("Versione in italiano semplice", [
    "Riscrittura dei contenuti in linguaggio piano (livello A2).",
    "Per parlanti italiano come seconda lingua e per la lettura facilitata.",
    "Sempre affiancata all'articolo completo.",
], N())

content("Contenuti in Lingua dei Segni (LIS)", [
    "Glossari e contenuti in Lingua dei Segni Italiana.",
    "Un canale di accesso dedicato alle persone sorde.",
    "Termini chiave del rischio spiegati in forma visiva.",
], N())

content("Versione in Braille", [
    "Versione stampabile in Braille degli articoli.",
    "Compatibile con le stampanti braille (formato BRF).",
    "Un canale dedicato a persone cieche e ipovedenti.",
], N())

content("Versioni in più lingue", [
    "Contenuti chiave tradotti in più lingue.",
    "Per residenti stranieri e visitatori.",
    "Lingua dichiarata correttamente anche per gli assistenti vocali.",
], N())

divider("Volontariato", N())

content("Diventa volontario", [
    "Come unirsi al Gruppo Comunale.",
    "Attività, formazione e impegno richiesto.",
    "Il valore del volontariato per la comunità.",
], N())

content("Affiliazioni e riconoscimenti europei", [
    "Quality Label del Corpo europeo di solidarietà (European Solidarity Corps) — Commissione europea.",
    "Codice di accreditamento dell'organizzazione: E10435833 (sempre accanto al logo, Reg. UE 2021/888).",
    "Servizio Nazionale di Protezione Civile (D.Lgs. 1/2018) — iscrizione al Registro Regionale del Volontariato.",
    "Coordinamento FEPIVOL — Federazione Pronto Intervento Volontariato ODV.",
], N(), fonti="European Solidarity Corps · DPC · Registro Regionale Lazio · FEPIVOL")

slide_sistema(N())

divider("Risorse e trasparenza", N())

content("Conoscere la Protezione Civile", [
    "Sezione dottrinale: storia, Servizio Nazionale e le quattro fasi (previsione, prevenzione, soccorso, superamento).",
    "La scienza del rischio e i Centri di competenza; la dimensione internazionale del sistema.",
    "Il catalogo dei rischi e il vulcanismo dei Colli Albani (Vulcano Laziale).",
    "Le telecomunicazioni in emergenza: capire come funziona la rete radio.",
], N(), fonti="DPC · INGV · ISPRA · CNR-IRPI · CIMA · CMCC · ISS")

content("FAQ e Glossario", [
    "Risposte alle domande frequenti dei cittadini.",
    "Glossario dei termini e delle sigle della Protezione Civile.",
    "Definizioni chiare, anche a comparsa sui termini tecnici nelle pagine.",
], N())

content("Normativa e Standard ISO", [
    "Raccolta della normativa nazionale, regionale e comunale.",
    "Schede degli standard internazionali ISO rilevanti per la materia.",
    "Dalla legge nazionale fino al livello comunale e internazionale.",
], N(), fonti="Normattiva · BURL Lazio · ISO")

content("Area download", [
    "Documenti ufficiali e modulistica scaricabili.",
    "Per ogni file sono indicati formato e dimensione (accessibilità).",
    "Un punto unico per atti, piani e materiali del Gruppo.",
], N())

content("Strumenti in tempo reale", [
    "Hub dei servizi di monitoraggio del territorio.",
    "Meteo, sismico, idrogeologico, incendi, qualità dell'aria, viabilità.",
    "Collegamenti diretti alle fonti ufficiali aggiornate.",
], N())

content("Audio e podcast", [
    "Contenuti ascoltabili per chi preferisce l'ascolto.",
    "Maggiore accessibilità all'informazione.",
    "Utili per anziani, ipovedenti e chi è in movimento.",
], N())

content("Open data e stato del sito", [
    "Dati aperti riutilizzabili dalla cittadinanza.",
    "Pagina di trasparenza sullo stato tecnico del sito.",
    "Riuso libero, secondo il principio di trasparenza della PA.",
], N())

content("Pagine legali e trasparenza", [
    "Privacy, note legali, dichiarazione di accessibilità, social media policy.",
    "Data di revisione sempre indicata.",
    "Conformità agli obblighi di trasparenza della Pubblica Amministrazione.",
], N())

divider("Comunicazione", N())

content("Comunicazioni e notizie", [
    "Sezione informativa costantemente aggiornata.",
    "Distinzione chiara tra informazione, prevenzione, allerta ed emergenza.",
    "Ogni avviso rimanda alla fonte ufficiale di riferimento.",
], N())

content("Comunicazione di crisi e social", [
    "Social media policy pubblica.",
    "Struttura dei messaggi secondo standard internazionali.",
    "Contrasto alla disinformazione con fonti ufficiali.",
], N(), fonti="ISO 22329 · CWA CEN/CENELEC · EENA (112)")

content("Canali di aggiornamento", [
    "Feed RSS per ogni sezione.",
    "Codici QR sugli articoli per l'accesso rapido.",
    "Condivisione semplice e senza tracciamento.",
], N())

divider("Infrastruttura, fonti e sicurezza", N())

content("Infrastruttura, sicurezza e prestazioni", [
    "Sito statico ad alte prestazioni, servito in HTTPS.",
    "Nessun cookie di profilazione: approccio privacy-first.",
    "Pubblicazione su doppio ambiente per la continuità del servizio.",
    "Contenuti documentabili anche in caso di audit esterno.",
], N())

content("Fonti e dati di riferimento", [
    "Istituzioni: DPC, INGV, ISPRA, CNR, ItaliaMeteo, EUMETSAT, Copernicus, ARPA Lazio, ASL Roma 6.",
    "Standard: WCAG 2.2; ISO 22320, 22324, 22329, 31000, 7010, 14090.",
    "Normativa: D.Lgs. 1/2018, Direttiva PCM 30/4/2021, Circolare DPC 6/8/2018, Legge Stanca.",
    "Design e linguaggio: AgID, Designers Italia, Bootstrap Italia.",
], N())

content("Punti di forza in sintesi", [
    "Accessibilità reale e multicanale: vista, udito, lingua, cognizione.",
    "Dati sempre verificati e con fonte dichiarata.",
    "Coerenza con le prassi ufficiali della Pubblica Amministrazione.",
    "Uno strumento di servizio pubblico, sobrio e affidabile.",
], N())

# chiusura
s = prs.slides.add_slide(BLANK); _norm(s)
rect(s, 0, 0, SW, SH, PRIMARY)
rect(s, 0, 0, Inches(0.22), SH, ACCENT)
s.shapes.add_picture(LOGO, Inches(1.0), Inches(0.7), height=Inches(1.15))
tf = textbox(s, Inches(1.0), Inches(2.2), Inches(11.3), Inches(1.2))
setrun(tf.paragraphs[0].add_run(), "Grazie per l'attenzione", 40, WHITE, bold=True)
rect(s, Inches(1.05), Inches(3.45), Inches(3.2), Pt(3), ACCENT)
# contatti reali (colonna sinistra)
tf2 = textbox(s, Inches(1.0), Inches(3.95), Inches(8.2), Inches(2.7))
p = tf2.paragraphs[0]
setrun(p.add_run(), "Gruppo Comunale Volontari di Protezione Civile di Genzano di Roma", 16, WHITE, bold=True)
for line in [
    "Sede: Via Sicilia, 13-15 — 00045 Genzano di Roma (RM)",
    "Segreteria: +39 06 9362600",
    "Email: segreteria@protezionecivilegenzano.it",
    "www.protezionecivilegenzano.it",
]:
    pp = tf2.add_paragraph(); pp.space_before = Pt(6); pp.line_spacing = 1.05
    setrun(pp.add_run(), line, 14, RGBColor(0xCF, 0xDD, 0xEC))
# riquadro 112 (colonna destra)
rect(s, Inches(9.9), Inches(3.95), Inches(2.6), Inches(2.0), ACCENT)
tb = textbox(s, Inches(9.9), Inches(3.95), Inches(2.6), Inches(2.0), anchor=MSO_ANCHOR.MIDDLE)
q = tb.paragraphs[0]; q.alignment = PP_ALIGN.CENTER
setrun(q.add_run(), "112", 46, PRIMARY_D, bold=True)
q2 = tb.add_paragraph(); q2.alignment = PP_ALIGN.CENTER
setrun(q2.add_run(), "In emergenza", 13, PRIMARY_D, bold=True)

# popola l'Agenda con PULSANTI-rettangolo cliccabili alle sezioni
def agenda_button(l, t, w, h, num, label, target):
    b = agenda.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    b.fill.solid(); b.fill.fore_color.rgb = PRIMARY
    b.line.color.rgb = ACCENT; b.line.width = Pt(1.25)
    b.shadow.inherit = False
    tfb = b.text_frame; tfb.word_wrap = True; tfb.vertical_anchor = MSO_ANCHOR.MIDDLE
    tfb.margin_left = Inches(0.24); tfb.margin_right = Inches(0.15)
    tfb.margin_top = Inches(0.02); tfb.margin_bottom = Inches(0.02)
    p = tfb.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r1 = p.add_run(); setrun(r1, f"{num:02d}   ", 16, ACCENT, bold=True)
    r2 = p.add_run(); setrun(r2, label, 15, WHITE, bold=True)
    b.click_action.target_slide = target
    add_shadow(b)

_cols = [Inches(0.85), Inches(6.95)]
_cw = Inches(5.55); _h = Inches(0.86); _y0 = 2.18; _step = 1.06
for _i, (_chap, _sl, _num) in enumerate(sec_targets):
    _col = 0 if _i < 4 else 1
    _row = _i if _i < 4 else _i - 4
    agenda_button(_cols[_col], Inches(_y0 + _row * _step), _cw, _h, _i + 1, _chap, _sl)
# riquadro vuoto di bilanciamento (colonna destra, 4ª riga): griglia 4+4 simmetrica
_fb = agenda.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                              _cols[1], Inches(_y0 + 3 * _step), _cw, _h)
_fb.fill.solid(); _fb.fill.fore_color.rgb = RGBColor(0xE7, 0xF0, 0xFA)
_fb.line.color.rgb = RGBColor(0xCB, 0xD7, 0xE6); _fb.line.width = Pt(1.0)
_fb.shadow.inherit = False

# ── Navigazione avanzata: menu di sezione sui divisori + frecce prec/succ ──
_all = list(prs.slides)
_dark_ids = {id(sl) for (_c, sl, _n) in sec_targets}  # divisori
_dark_ids.add(id(_all[-1]))                            # + chiusura


def _menu_chip(slide, x, y, w, h, idx, label, target, current):
    b = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    b.shadow.inherit = False
    tf = b.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.16); tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.01); tf.margin_bottom = Inches(0.01)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    if current:
        b.fill.solid(); b.fill.fore_color.rgb = ACCENT; b.line.fill.background()
        r1 = p.add_run(); setrun(r1, f"{idx:02d}  ", 11, PRIMARY_D, bold=True)
        r2 = p.add_run(); setrun(r2, label, 11, PRIMARY_D, bold=True)
    else:
        b.fill.solid(); b.fill.fore_color.rgb = PRIMARY
        b.line.color.rgb = RGBColor(0x4A, 0x6A, 0x8A); b.line.width = Pt(0.75)
        r1 = p.add_run(); setrun(r1, f"{idx:02d}  ", 11, ACCENT, bold=True)
        r2 = p.add_run(); setrun(r2, label, 11, WHITE)
        b.click_action.target_slide = target


# menu di tutte le sezioni su ogni divisore (corrente evidenziata)
for (_c, _dsl, _n) in sec_targets:
    for j, (cj, sj, nj) in enumerate(sec_targets):
        _menu_chip(_dsl, Inches(7.05), Inches(1.62 + j * 0.6), Inches(5.45),
                   Inches(0.5), j + 1, cj, sj, sj is _dsl)

# frecce precedente / successiva su tutte le slide (tranne copertina)
for i, sl in enumerate(_all):
    if i == 0:
        continue
    col = ACCENT if id(sl) in _dark_ids else PRIMARY
    tb = textbox(sl, Inches(10.78), Inches(6.97), Inches(0.36), Inches(0.4))
    pp = tb.paragraphs[0]; pp.alignment = PP_ALIGN.CENTER
    rr = pp.add_run(); setrun(rr, "‹", 15, col, bold=True)
    link_to_slide(rr, sl, _all[i - 1])
    if i < len(_all) - 1:
        tb2 = textbox(sl, Inches(11.18), Inches(6.97), Inches(0.36), Inches(0.4))
        pp2 = tb2.paragraphs[0]; pp2.alignment = PP_ALIGN.CENTER
        rr2 = pp2.add_run(); setrun(rr2, "›", 15, col, bold=True)
        link_to_slide(rr2, sl, _all[i + 1])

out = os.path.join(REPO, "static/manuali/presentazione-struttura-sito.pptx")
os.makedirs(os.path.dirname(out), exist_ok=True)
prs.save(out)
print("OK", out, "·", len(prs.slides.__iter__.__self__._sldIdLst), "slide")

import shutil, subprocess
_pdf = out[:-5] + ".pdf"
_soffice = shutil.which("libreoffice") or shutil.which("soffice")
if _soffice:
    subprocess.run([_soffice, "--headless", "--convert-to", "pdf", "--outdir",
                    os.path.dirname(out), out],
                   stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
    print("PDF:", _pdf if os.path.exists(_pdf) else "(conversione PDF non riuscita)")
    # Allinea automaticamente la dimensione mostrata nell'Area Download al peso
    # reale del PDF, così l'etichetta non si disallinea più a ogni rigenerazione.
    if os.path.exists(_pdf):
        _label = ("%.1f" % (os.path.getsize(_pdf) / 1048576)).replace(".", ",")
        _adl = os.path.join(REPO, "content/area-download/_index.md")
        try:
            _txt = open(_adl, encoding="utf-8").read()
            _new = re.sub(r"(presentazione-struttura-sito\.pdf\).*?PDF · )[0-9.,]+( MB)",
                          lambda m: m.group(1) + _label + m.group(2), _txt)
            if _new != _txt:
                open(_adl, "w", encoding="utf-8").write(_new)
                print("Area download: dimensione PDF allineata →", _label, "MB")
        except OSError as _e:
            print("Avviso: dimensione in area-download non aggiornata:", _e)
else:
    print("LibreOffice assente: converti manualmente il .pptx in PDF.")
